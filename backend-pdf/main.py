import openpyxl
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import fitz
import img2pdf
import io
import os
import tempfile
import zipfile
from PIL import Image
from typing import List

app = FastAPI(title="PDF Service", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def root():
    return {"status": "ok", "service": "PDF Backend"}

@app.get("/health")
def health():
    return {"status": "healthy"}

@app.post("/merge")
async def merge_pdf(files: List[UploadFile] = File(...)):
    merged = fitz.open()
    for file in files:
        data = await file.read()
        doc = fitz.open(stream=data, filetype="pdf")
        merged.insert_pdf(doc)
        doc.close()
    buf = io.BytesIO()
    merged.save(buf)
    merged.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=merged.pdf"})

@app.post("/split")
async def split_pdf(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, 'w') as zf:
        for i, page in enumerate(doc):
            single = fitz.open()
            single.insert_pdf(doc, from_page=i, to_page=i)
            page_buf = io.BytesIO()
            single.save(page_buf)
            single.close()
            zf.writestr(f"page_{i+1}.pdf", page_buf.getvalue())
    doc.close()
    zip_buf.seek(0)
    return StreamingResponse(zip_buf, media_type="application/zip",
        headers={"Content-Disposition": "attachment; filename=pages.zip"})

@app.post("/compress")
async def compress_pdf(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    buf = io.BytesIO()
    doc.save(buf, garbage=4, deflate=True, clean=True)
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=compressed.pdf"})

@app.post("/pdf-to-image")
async def pdf_to_image(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, 'w') as zf:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i+1}.png", pix.tobytes("png"))
    doc.close()
    zip_buf.seek(0)
    return StreamingResponse(zip_buf, media_type="application/zip",
        headers={"Content-Disposition": "attachment; filename=images.zip"})

@app.post("/image-to-pdf")
async def image_to_pdf(files: List[UploadFile] = File(...)):
    images = []
    for file in files:
        data = await file.read()
        img = Image.open(io.BytesIO(data)).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=95)
        images.append(buf.getvalue())
    pdf_bytes = img2pdf.convert(images)
    return StreamingResponse(io.BytesIO(pdf_bytes), media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=converted.pdf"})

@app.post("/rotate")
async def rotate_pdf(files: List[UploadFile] = File(...), degrees: int = Form(90)):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    for page in doc:
        page.set_rotation(page.rotation + degrees)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=rotated.pdf"})

@app.post("/pdf-to-word")
async def pdf_to_word(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(data)
        tmp_pdf_path = tmp_pdf.name
    tmp_docx_path = tmp_pdf_path.replace(".pdf", ".docx")
    try:
        from pdf2docx import Converter
        cv = Converter(tmp_pdf_path)
        cv.convert(tmp_docx_path)
        cv.close()
        with open(tmp_docx_path, "rb") as f:
            docx_bytes = f.read()
    finally:
        os.unlink(tmp_pdf_path)
        if os.path.exists(tmp_docx_path):
            os.unlink(tmp_docx_path)
    return StreamingResponse(io.BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=converted.docx"})

@app.post("/word-to-pdf")
async def word_to_pdf(files: List[UploadFile] = File(...)):
    from docx import Document
    data = await files[0].read()
    doc_word = Document(io.BytesIO(data))
    full_text = [para.text for para in doc_word.paragraphs]
    for table in doc_word.tables:
        for row in table.rows:
            full_text.append(" | ".join(c.text.strip() for c in row.cells))
    pdf_doc = fitz.open()
    page = pdf_doc.new_page(width=595, height=842)
    y, margin = 50, 50
    for line in full_text:
        if not line.strip():
            y += 9
            continue
        page.insert_text((margin, y), line[:100], fontsize=11, color=(0, 0, 0))
        y += 18
        if y > 800:
            page = pdf_doc.new_page(width=595, height=842)
            y = 50
    buf = io.BytesIO()
    pdf_doc.save(buf)
    pdf_doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=converted.pdf"})

@app.post("/pdf-to-excel")
async def pdf_to_excel(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "PDF Data"
    ws.append(["Page", "Line", "Content"])
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        for line_num, line in enumerate(lines, start=1):
            ws.append([page_num, line_num, line])
    doc.close()
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=extracted.xlsx"})

@app.post("/excel-to-pdf")
async def excel_to_pdf(files: List[UploadFile] = File(...)):
    data = await files[0].read()
    wb = openpyxl.load_workbook(io.BytesIO(data))
    pdf_doc = fitz.open()
    page_width, page_height = 842, 595
    margin = 30
    for sheet in wb.worksheets:
        page = pdf_doc.new_page(width=page_width, height=page_height)
        y = margin + 20
        page.insert_text((margin, margin), f"Sheet: {sheet.title}",
                         fontsize=13, color=(0.1, 0.1, 0.6))
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        col_count = max(len(r) for r in rows)
        col_width = min(120, (page_width - margin * 2) // max(col_count, 1))
        for row_idx, row in enumerate(rows):
            x = margin
            for col_idx, cell in enumerate(row):
                cell_text = str(cell) if cell is not None else ""
                cell_text = cell_text[:18]
                color = (0.1, 0.1, 0.6) if row_idx == 0 else (0.15, 0.15, 0.15)
                page.insert_text((x + 3, y), cell_text, fontsize=9, color=color)
                page.draw_rect(fitz.Rect(x, y - 12, x + col_width, y + 4),
                               color=(0.8, 0.8, 0.8), width=0.3)
                x += col_width
            y += 18
            if y > page_height - margin:
                page = pdf_doc.new_page(width=page_width, height=page_height)
                y = margin + 20
    buf = io.BytesIO()
    pdf_doc.save(buf)
    pdf_doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=converted.pdf"})

@app.post("/watermark")
async def watermark_pdf(
    files: List[UploadFile] = File(...),
    watermark_text: str = Form("CONFIDENTIAL")
):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    for page in doc:
        w, h = page.rect.width, page.rect.height
        page.insert_text(
            (w * 0.1, h * 0.5), watermark_text,
            fontsize=60, color=(0.8, 0.8, 0.8),
            rotate=45, overlay=True
        )
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=watermarked.pdf"})

@app.post("/page-numbers")
async def add_page_numbers(
    files: List[UploadFile] = File(...),
    position: str = Form("bottom-center")
):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    total = len(doc)
    for i, page in enumerate(doc):
        w, h = page.rect.width, page.rect.height
        text = f"{i+1} / {total}"
        if position == "bottom-center":
            x, y = w / 2 - 20, h - 20
        elif position == "bottom-right":
            x, y = w - 60, h - 20
        else:
            x, y = 20, h - 20
        page.insert_text((x, y), text, fontsize=11, color=(0.3, 0.3, 0.3))
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=numbered.pdf"})

@app.post("/protect")
async def protect_pdf(
    files: List[UploadFile] = File(...),
    password: str = Form("1234")
):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    buf = io.BytesIO()
    doc.save(buf, encryption=fitz.PDF_ENCRYPT_AES_256,
             user_pw=password, owner_pw=password + "_owner")
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=protected.pdf"})

@app.post("/unlock")
async def unlock_pdf(
    files: List[UploadFile] = File(...),
    password: str = Form("")
):
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    if doc.is_encrypted:
        if not doc.authenticate(password):
            return {"error": "Wrong password"}
    buf = io.BytesIO()
    doc.save(buf, encryption=fitz.PDF_ENCRYPT_NONE)
    doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=unlocked.pdf"})

@app.post("/pdf-to-ppt")
async def pdf_to_ppt(files: List[UploadFile] = File(...)):
    from pptx import Presentation
    from pptx.util import Inches
    data = await files[0].read()
    doc = fitz.open(stream=data, filetype="pdf")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for page in doc:
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img_buf = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_buf, 0, 0,
            width=prs.slide_width, height=prs.slide_height)
    doc.close()
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=converted.pptx"})

@app.post("/ppt-to-pdf")
async def ppt_to_pdf(files: List[UploadFile] = File(...)):
    from pptx import Presentation
    data = await files[0].read()
    prs = Presentation(io.BytesIO(data))
    pdf_doc = fitz.open()
    page_w = int(prs.slide_width.pt)
    page_h = int(prs.slide_height.pt)
    for slide in prs.slides:
        page = pdf_doc.new_page(width=page_w, height=page_h)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    fontsize = 16
                    try:
                        if para.runs:
                            fs = para.runs[0].font.size
                            if fs:
                                fontsize = min(int(fs.pt), 48)
                    except:
                        pass
                    x = int(shape.left.pt) if shape.left else 20
                    y = int(shape.top.pt) + 20 if shape.top else 40
                    page.insert_text((x, y), text, fontsize=fontsize, color=(0, 0, 0))
    buf = io.BytesIO()
    pdf_doc.save(buf)
    pdf_doc.close()
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=converted.pdf"})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
